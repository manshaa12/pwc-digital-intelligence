from pptx import Presentation
from pptx.util import Inches
import pandas as pd

# Create a presentation object
prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Analysis of Adverse Effects of Tramal and Lyrica"
subtitle.text = "Quantitative Finance Team\nDate: July 15, 2024"

# Slide 1: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Introduction"
content = slide.placeholders[1]
content.text = ("- Objective: Analyze and compare adverse effects of Tramal and Lyrica using FAERS data from 2019\n"
                "- Steps undertaken for the analysis:\n"
                "  1. Descriptive overview of Tramal adverse effects\n"
                "  2. Comparison with Lyrica\n"
                "  3. Suggested further investigations")

# Slide 2: Adverse Effects of Tramal
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Adverse Effects of Tramal"
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
text_frame = content.text_frame
text_frame.text = "Top 10 Adverse Effects of Tramal in 2019"

# Add the bar plot image for Tramal
slide.shapes.add_picture("tramal_adverse_effects.png", Inches(0.5), Inches(2.5), Inches(9), Inches(4))

# Slide 3: Adverse Effects of Lyrica
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Adverse Effects of Lyrica"
content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
text_frame = content.text_frame
text_frame.text = "Top 10 Adverse Effects of Lyrica in 2019"

# Add the bar plot image for Lyrica
slide.shapes.add_picture("lyrica_adverse_effects.png", Inches(0.5), Inches(2.5), Inches(9), Inches(4))

# Slide 4: Comparison of Adverse Effects
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Comparison of Adverse Effects"
content = slide.placeholders[1]
content.text = ("- Comparison of top adverse effects of Tramal and Lyrica\n"
                "- Table or chart comparing the top adverse effects\n"
                "- Key insights from the comparison")

# Slide 5: Further Investigations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Further Investigations"
content = slide.placeholders[1]
content.text = ("- Suggested further investigations:\n"
                "  1. Longitudinal studies\n"
                "  2. Subgroup analysis\n"
                "  3. Dose-response relationship\n"
                "  4. Comparative effectiveness research\n"
                "  5. Real-world evidence\n"
                "- Importance of each suggested investigation")

# Save the presentation
prs.save('/mnt/data/Adverse_Effects_Analysis_Presentation.pptx')
